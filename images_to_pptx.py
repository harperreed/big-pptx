
from pptx import Presentation
from natsort import natsorted, ns
from pptx.util import Pt, Inches
import pptx
import pptx.util
import glob
import scipy.misc

from slugify import slugify

import argparse
import logging


try:
    import argparse
    parser = argparse.ArgumentParser( description='Harper Rules, LLC Presentation Generator!')
except ImportError:
    flags = None

loglevel = logging.INFO

logging.basicConfig(level=loglevel)
logger = logging.getLogger(__name__)

class presentation_generator:

    def __init__(self, presentation_name, slides_path, root_path="./", slide_pattern="slide*", width=1920, height=1080):
        self.presentation_name = presentation_name
        self.slug = slugify(unicode(self.presentation_name))
        self.width = width
        self.height = height
        self.slides_path = slides_path
        self.root_path = root_path
        self.slide_pattern = slide_pattern

    def build(self):
        # new
        logging.info("Generating PPTX")
        prs = pptx.Presentation()
 
        prs.slide_height = Pt(self.height)
        prs.slide_width = Pt(self.width)


        logging.info("Grabbing slides")
        slides = natsorted(glob.glob(self.slides_path + self.slide_pattern))


        logging.info("Adding slides")
        for slide_image in slides:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            pic   = slide.shapes.add_picture(slide_image, Pt(0), Pt(0))
            print "added slide"


        filename = self.root_path + self.slug + "-" + str(self.width) + "x" + str(self.height) + ".pptx"
        logging.info("Saving PPTX: " + filename)
        prs.save(filename)
        


def main():

    parser.add_argument("-r","--root_path", help="Path of the presentation", required=True)
    parser.add_argument("-s","--slides_path", help="Path of the slide images", required=True)
    parser.add_argument("-n","--name", help="Name of the presentation", required=True)
    flags = parser.parse_args()

    name = flags.name
    root_path = flags.root_path
    slides_path = flags.slides_path

    presentation = presentation_generator(presentation_name=name, root_path=root_path, slides_path=slides_path )

    presentation.build()


if __name__ == '__main__':
    main()
